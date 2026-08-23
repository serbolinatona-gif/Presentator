"""
pptx_builder.py
Сборка .pptx с двумя независимыми слоями:
1. Стиль (style) — задаёт фон, палитру, шрифты и декоративный "фрейм" слайда
   (draw_style_frame). Раньше стили были почти незаметны, потому что фоновая
   картинка на весь слайд перекрывала всё оформление — теперь картинка ставится
   в отдельную область рядом с текстом, а не на фон.
2. Раскладка (layout) — задаёт геометрию/структуру конкретного слайда (9 типов:
   title, toc, thesis_proof, dense_text, concept_anatomy, comparison,
   causal_chain, quote_context, conclusion). Раскладка использует цвета/шрифты
   стиля, но сама геометрия одинакова независимо от стиля — так 5 стилей x 9
   раскладок не превращаются в 45 отдельных функций.

Известное упрощение: для concept_anatomy не рисуются линии-выноски (leader lines)
от картинки к пунктам — в python-pptx это делается через произвольные connector-
shapes с точными координатами и на практике выглядит хрупко при разных
пропорциях картинок. Вместо этого части располагаются явным аккуратным списком
под определением, что тоже читаемо и профессионально выглядит.
"""

import io
import logging
import math
import re
from typing import List, Optional

import httpx
from lxml import etree
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

logger = logging.getLogger("slideforge.pptx_builder")

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.7)

# ---------------------------------------------------------------------------
# Стили — палитра/шрифты/декор
# ---------------------------------------------------------------------------

STYLE_TOKENS = {
    "minimal": {
        "bg": "F9F9F9", "bg2": None, "title": "1A1A1A", "text": "333333",
        "accent": "1D4ED8", "muted": "9CA3AF", "on_accent": "FFFFFF",
        "font_title": "Helvetica", "font_body": "Helvetica",
        "mode": "light", "decoration": "none",
    },
    "academic": {
        "bg": "FAF9F6", "bg2": None, "title": "1D3557", "text": "222222",
        "accent": "1D3557", "muted": "6B7280", "on_accent": "FFFFFF",
        "font_title": "Georgia", "font_body": "Georgia",
        "mode": "light", "decoration": "rules",
    },
    "creative": {
        "bg": "FFFFFF", "bg2": None, "title": "18122B", "text": "3D3355",
        "accent": "6C5CE7", "muted": "9B93B8", "on_accent": "FFFFFF",
        "font_title": "Verdana", "font_body": "Verdana",
        "mode": "light", "decoration": "geometry",
        "palette": ["FF6B6B", "4ECDC4", "FFD93D", "6C5CE7", "1DD3B0", "FF9F43", "3AB0FF", "F368E0"],
    },
    "corporate": {
        "bg": "FFFFFF", "bg2": None, "title": "0F172A", "text": "1E293B",
        "accent": "0EA5E9", "muted": "64748B", "on_accent": "FFFFFF",
        "font_title": "Calibri", "font_body": "Calibri",
        "mode": "light", "decoration": "cards",
    },
    "dark": {
        "bg": "121212", "bg2": None, "title": "FFFFFF", "text": "E0E0E0",
        "accent": "22D3EE", "muted": "9CA3AF", "on_accent": "111111",
        "font_title": "Calibri", "font_body": "Calibri",
        "mode": "dark", "decoration": "glow",
    },
}


def _clean_text(text) -> str:
    """
    Страховка от типографских багов в тексте от AI:
    - схлопывает переносы строк/лишние пробелы в один пробел (чтобы не было
      "склеенных" абзацев без пробела при копировании/повторной генерации)
    - делает первую букву предложения заглавной, если модель вернула её строчной
    """
    if not text:
        return ""
    s = str(text).replace("\r", " ").replace("\n", " ")
    s = re.sub(r"\s+", " ", s).strip()
    if s and s[0].isalpha() and s[0].islower():
        s = s[0].upper() + s[1:]
    return s


def _fit_font_size(char_count: int, base_size: int, width_in: float, height_in: float, min_size: int = 10) -> int:
    """
    Расчётное (не полагающееся на автоподгонку PowerPoint при открытии) уменьшение
    размера шрифта, если текста больше, чем помещается в рамку. PowerPoint autofit
    работает только при открытии файла и не везде одинаково (Google Slides/LibreOffice
    иногда игнорируют флаг) — поэтому считаем сами, по приблизительной модели
    "символов на строку" и "строк в высоту".
    """
    if char_count <= 0:
        return base_size
    avg_char_width_factor = 0.52  # эмпирический коэффициент средней ширины символа к кеглю
    line_height_factor = 1.35
    width_pt = width_in * 72
    height_pt = height_in * 72
    chars_per_line = max(1, width_pt / (base_size * avg_char_width_factor))
    lines_needed = max(1, math.ceil(char_count / chars_per_line))
    max_lines_fit = max(1, int(height_pt / (base_size * line_height_factor)))
    if lines_needed <= max_lines_fit:
        return base_size
    scale = math.sqrt(max_lines_fit / lines_needed)
    return max(min_size, int(base_size * scale))


def _estimate_block_height_in(char_count: int, font_size: int, width_in: float, line_height_factor: float = 1.35) -> float:
    """Оценка высоты (в дюймах), которую займёт текст такой длины при данном кегле и ширине."""
    if char_count <= 0:
        return 0.0
    avg_char_width_factor = 0.52
    width_pt = width_in * 72
    chars_per_line = max(1, width_pt / (font_size * avg_char_width_factor))
    lines = max(1, math.ceil(char_count / chars_per_line))
    return (lines * font_size * line_height_factor) / 72


def _hex(color_hex: str) -> RGBColor:
    return RGBColor.from_string(color_hex.lstrip("#"))


EMU_PER_INCH = 914400


def _to_inches(value) -> float:
    """
    Устойчиво конвертирует в дюймы. Length-объекты python-pptx (Inches/Pt/Emu) теряют
    свой тип и становятся обычным int (EMU) после арифметики (например Inches(5) - Inches(1)
    возвращает int, а не Length) — .inches у такого int уже не будет. Поэтому здесь
    используем .inches, если он есть, иначе делим сырые EMU вручную.
    """
    if hasattr(value, "inches"):
        return value.inches
    return value / EMU_PER_INCH


def _set_shape_transparency(shape, alpha_percent: float):
    """alpha_percent: 0.0 (прозрачно) .. 1.0 (непрозрачно). Нет публичного API в python-pptx."""
    sp = shape.fill.fore_color._xFill
    srgb = sp.find("{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr")
    if srgb is None:
        return
    alpha_el = etree.SubElement(srgb, "{http://schemas.openxmlformats.org/drawingml/2006/main}alpha")
    alpha_el.set("val", str(int(alpha_percent * 100000)))


async def _download_image(url: str) -> Optional[bytes]:
    try:
        async with httpx.AsyncClient(timeout=20.0) as client:
            resp = await client.get(url, follow_redirects=True)
            if resp.status_code == 200:
                return resp.content
    except httpx.RequestError as exc:
        logger.warning("Не удалось скачать изображение %s: %s", url, exc)
    return None


# ---------------------------------------------------------------------------
# Базовые примитивы вёрстки
# ---------------------------------------------------------------------------

def _solid_bg(slide, color_hex):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _hex(color_hex)


def _gradient_bg(slide, c1, c2, angle=45.0):
    fill = slide.background.fill
    fill.gradient()
    stops = fill.gradient_stops
    stops[0].color.rgb = _hex(c1)
    stops[1].color.rgb = _hex(c2)
    fill.gradient_angle = angle


def _text(slide, text, left, top, width, height, size, color, font, bold=False, italic=False,
          align=PP_ALIGN.LEFT, line_spacing=1.15):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE  # размер считаем сами — надёжнее автоподгонки PowerPoint,
                                        # которую не все программы (Google Slides/LibreOffice) применяют одинаково
    clean = _clean_text(text)
    fitted_size = _fit_font_size(len(clean), size, _to_inches(width), _to_inches(height))
    p = tf.paragraphs[0]
    p.text = clean
    p.font.size = Pt(fitted_size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font
    p.font.color.rgb = _hex(color)
    p.alignment = align
    p.line_spacing = line_spacing
    box.text_frame.word_wrap = True
    return box, _estimate_block_height_in(len(clean), fitted_size, _to_inches(width), line_spacing)


def _multi_para(slide, lines, left, top, width, height, size, color, font, bold=False,
                 marker="", space_after=10, align=PP_ALIGN.LEFT):
    if isinstance(lines, str):
        lines = [lines]  # защита: если AI вернул строку вместо списка, не разбиваем по буквам
    clean_lines = [f"{marker}{_clean_text(line)}" for line in lines]
    total_chars = sum(len(l) for l in clean_lines)
    fitted_size = _fit_font_size(total_chars, size, _to_inches(width), _to_inches(height))

    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    for i, line in enumerate(clean_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(fitted_size)
        p.font.bold = bold
        p.font.name = font
        p.font.color.rgb = _hex(color)
        p.space_after = Pt(space_after)
        p.alignment = align
    est_height = _estimate_block_height_in(total_chars, fitted_size, _to_inches(width)) + (len(clean_lines) * space_after / 72)
    return box, est_height


def _rect(slide, left, top, width, height, color_hex, transparency=None, line=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = _hex(color_hex)
    if transparency is not None:
        _set_shape_transparency(shp, transparency)
    if not line:
        shp.line.fill.background()
    return shp


def _oval(slide, left, top, width, height, color_hex, transparency=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = _hex(color_hex)
    shp.line.fill.background()
    if transparency is not None:
        _set_shape_transparency(shp, transparency)
    return shp


def _image_cover(slide, image_bytes, left, top, width, height):
    """Вставляет картинку с заполнением области без искажения пропорций (crop-to-fill)."""
    pic = slide.shapes.add_picture(io.BytesIO(image_bytes), left, top, width=width, height=height)
    try:
        img = Image.open(io.BytesIO(image_bytes))
        iw, ih = img.size
        box_ratio = width / height
        img_ratio = iw / ih
        if img_ratio > box_ratio:
            crop = (1 - box_ratio / img_ratio) / 2
            pic.crop_left = crop
            pic.crop_right = crop
        else:
            crop = (1 - img_ratio / box_ratio) / 2
            pic.crop_top = crop
            pic.crop_bottom = crop
    except Exception:
        logger.warning("Не удалось вычислить пропорции картинки для crop-to-fill")
    return pic


def _image_or_placeholder(slide, image_bytes, left, top, width, height, tokens):
    if image_bytes:
        _image_cover(slide, image_bytes, left, top, width, height)
    else:
        accent2 = tokens.get("bg2") or tokens["accent"]
        _gradient_bg  # noqa: touch to avoid unused import pruning by linters
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shp.fill.gradient()
        stops = shp.fill.gradient_stops
        stops[0].color.rgb = _hex(tokens["accent"])
        stops[1].color.rgb = _hex(accent2)
        shp.line.fill.background()


# ---------------------------------------------------------------------------
# Фрейм стиля — фон + декор, вызывается для КАЖДОГО слайда до раскладки
# ---------------------------------------------------------------------------

def draw_style_frame(slide, style: str, tokens: dict, idx: int, slide_count: int, is_title: bool):
    if tokens["bg2"] and tokens["decoration"] == "geometry":
        _gradient_bg(slide, tokens["bg"], tokens["bg2"], angle=45.0 if idx % 2 == 0 else 135.0)
    else:
        _solid_bg(slide, tokens["bg"])

    deco = tokens["decoration"]

    if deco == "rules":  # academic
        _rect(slide, 0, 0, SLIDE_W, Pt(3), tokens["accent"])
        if not is_title:
            _text(slide, f"{idx + 1:02d} / {slide_count:02d}", SLIDE_W - Inches(1.4), SLIDE_H - Inches(0.45),
                  Inches(1.1), Inches(0.35), 10, tokens["muted"], tokens["font_body"], align=PP_ALIGN.RIGHT)
            _rect(slide, MARGIN, SLIDE_H - Inches(0.55), SLIDE_W - 2 * MARGIN, Pt(1), tokens["muted"])

    elif deco == "cards":  # corporate
        _rect(slide, 0, 0, SLIDE_W, Inches(0.12), tokens["accent"])
        if not is_title:
            _text(slide, f"{idx + 1:02d}", SLIDE_W - Inches(1.0), SLIDE_H - Inches(0.5),
                  Inches(0.7), Inches(0.35), 11, tokens["muted"], tokens["font_body"], align=PP_ALIGN.RIGHT)

    elif deco == "geometry":  # creative — светлый фон + разноцветные фигуры (Figma-style)
        palette = tokens.get("palette") or [tokens["accent"]]

        def pcolor(offset):
            return palette[(idx + offset) % len(palette)]

        # верхний правый угол — крупный полупрозрачный круг
        _oval(slide, SLIDE_W - Inches(2.6), Inches(-1.3), Inches(3.6), Inches(3.6), pcolor(0), 0.85)
        # маленький акцентный круг
        _oval(slide, SLIDE_W - Inches(1.5), Inches(1.6), Inches(0.9), Inches(0.9), pcolor(1), 0.9)
        # скруглённый прямоугольник снизу слева
        rrect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(-0.6), SLIDE_H - Inches(1.4),
                                        Inches(2.6), Inches(2.6))
        rrect.rotation = 20
        rrect.fill.solid()
        rrect.fill.fore_color.rgb = _hex(pcolor(2))
        rrect.line.fill.background()
        _set_shape_transparency(rrect, 0.8)
        # маленький треугольник-акцент
        tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(0.3), SLIDE_H - Inches(0.9),
                                      Inches(0.7), Inches(0.7))
        tri.rotation = -15
        tri.fill.solid()
        tri.fill.fore_color.rgb = _hex(pcolor(3))
        tri.line.fill.background()
        # тонкая цветная плашка-акцент под заголовком (если не титульный)
        if not is_title:
            _rect(slide, MARGIN, Inches(0.78), Inches(0.55), Inches(0.1), pcolor(0))

    elif deco == "glow":  # dark
        glow = _oval(slide, SLIDE_W - Inches(2.6), Inches(-1.8), Inches(4.0), Inches(4.0), tokens["accent"], 0.14)


def _content_top(tokens) -> Inches:
    if tokens["decoration"] == "rules":
        return Inches(1.05)
    if tokens["decoration"] == "cards":
        return Inches(1.0)
    return Inches(0.9)


def _content_bottom(tokens) -> Inches:
    if tokens["decoration"] == "rules":
        return SLIDE_H - Inches(0.75)
    return SLIDE_H - Inches(0.5)


# ---------------------------------------------------------------------------
# Раскладки
# ---------------------------------------------------------------------------

def _render_title(slide, data, tokens, image_bytes):
    _text(slide, data.get("title", ""), Inches(1.0), Inches(2.9), SLIDE_W - Inches(2.0), Inches(1.5),
          44, tokens["title"], tokens["font_title"], bold=True, align=PP_ALIGN.CENTER)
    subtitle = data.get("subtitle")
    if subtitle:
        _text(slide, subtitle, Inches(1.8), Inches(4.3), SLIDE_W - Inches(3.6), Inches(0.9),
              18, tokens["text"], tokens["font_body"], align=PP_ALIGN.CENTER)


def _render_toc(slide, data, tokens, image_bytes):
    top = _content_top(tokens)
    _text(slide, data.get("title", "Оглавление"), MARGIN, top, SLIDE_W - 2 * MARGIN, Inches(0.9),
          30, tokens["title"], tokens["font_title"], bold=True)
    items = data.get("items") or []
    y = top + Inches(1.1)
    row_h = Inches(0.62)
    for i, item in enumerate(items):
        _text(slide, f"{i + 1:02d}", MARGIN, y, Inches(0.7), row_h, 18, tokens["accent"], tokens["font_title"], bold=True)
        _text(slide, item, MARGIN + Inches(0.9), y, SLIDE_W - 2 * MARGIN - Inches(0.9), row_h,
              16, tokens["text"], tokens["font_body"])
        y += row_h


def _render_thesis_proof(slide, data, tokens, image_bytes):
    top = _content_top(tokens)
    left_w = Inches(5.7)
    claim_box_h = Inches(2.0)
    _, claim_used_h = _text(slide, data.get("claim", data.get("title", "")), MARGIN, top, left_w, claim_box_h,
                             26, tokens["title"], tokens["font_title"], bold=True, line_spacing=1.2)
    facts_top = top + Inches(max(claim_used_h + 0.25, 1.4))  # не ближе фиксированного минимума, но и не наезжаем
    facts = data.get("facts") or []
    _multi_para(slide, facts, MARGIN, facts_top, left_w, _content_bottom(tokens) - facts_top,
                16, tokens["text"], tokens["font_body"], marker="—  ", space_after=12)
    img_left = MARGIN + left_w + Inches(0.4)
    img_w = SLIDE_W - MARGIN - img_left
    _image_or_placeholder(slide, image_bytes, img_left, top, img_w, _content_bottom(tokens) - top, tokens)


def _render_dense_text(slide, data, tokens, image_bytes):
    top = _content_top(tokens)
    _text(slide, data.get("problem", data.get("title", "")), MARGIN, top, SLIDE_W - 2 * MARGIN, Inches(1.0),
          26, tokens["title"], tokens["font_title"], bold=True)
    body_top = top + Inches(1.15)
    body_bottom = _content_bottom(tokens)
    paragraph = data.get("paragraph", "")
    if image_bytes:
        text_w = Inches(7.2)
        _text(slide, paragraph, MARGIN, body_top, text_w, body_bottom - body_top,
              15, tokens["text"], tokens["font_body"], line_spacing=1.3)
        img_left = MARGIN + text_w + Inches(0.4)
        img_w = SLIDE_W - MARGIN - img_left
        _image_or_placeholder(slide, image_bytes, img_left, body_top, img_w, body_bottom - body_top, tokens)
    else:
        _text(slide, paragraph, MARGIN, body_top, SLIDE_W - 2 * MARGIN, body_bottom - body_top,
              16, tokens["text"], tokens["font_body"], line_spacing=1.35)


def _render_concept_anatomy(slide, data, tokens, image_bytes):
    top = _content_top(tokens)
    left_w = Inches(5.9)
    _text(slide, data.get("term", data.get("title", "")), MARGIN, top, left_w, Inches(0.6),
          22, tokens["title"], tokens["font_title"], bold=True)
    def_top = top + Inches(0.65)
    def_box_h = Inches(1.7)
    _, def_used_h = _text(slide, data.get("definition", ""), MARGIN, def_top, left_w, def_box_h,
                           15, tokens["text"], tokens["font_body"], line_spacing=1.25)

    parts = data.get("parts") or []
    py = def_top + Inches(max(def_used_h + 0.3, 1.0))
    for part in parts:
        name = part.get("name", "") if isinstance(part, dict) else str(part)
        func = part.get("function", "") if isinstance(part, dict) else ""
        _text(slide, name, MARGIN, py, left_w, Inches(0.3), 14, tokens["accent"], tokens["font_body"], bold=True)
        _text(slide, func, MARGIN, py + Inches(0.3), left_w, Inches(0.4), 13, tokens["text"], tokens["font_body"])
        py += Inches(0.75)

    img_left = MARGIN + left_w + Inches(0.4)
    img_w = SLIDE_W - MARGIN - img_left
    _image_or_placeholder(slide, image_bytes, img_left, top, img_w, _content_bottom(tokens) - top, tokens)


def _render_comparison(slide, data, tokens, image_bytes):
    top = _content_top(tokens)
    _text(slide, data.get("summary", data.get("title", "")), MARGIN, top, SLIDE_W - 2 * MARGIN, Inches(0.9),
          24, tokens["title"], tokens["font_title"], bold=True)

    cols_top = top + Inches(1.0)
    img_h = Inches(1.5) if image_bytes else Inches(0)
    if image_bytes:
        _image_or_placeholder(slide, image_bytes, MARGIN, cols_top, SLIDE_W - 2 * MARGIN, img_h, tokens)
        cols_top += img_h + Inches(0.25)

    col_w = (SLIDE_W - 2 * MARGIN - Inches(0.6)) / 2
    right_left = MARGIN + col_w + Inches(0.6)
    bottom = _content_bottom(tokens)

    _text(slide, data.get("left_label", ""), MARGIN, cols_top, col_w, Inches(0.4),
          17, tokens["accent"], tokens["font_title"], bold=True)
    _multi_para(slide, data.get("left_points") or [], MARGIN, cols_top + Inches(0.45), col_w, bottom - cols_top - Inches(0.45),
                14, tokens["text"], tokens["font_body"], marker="•  ", space_after=8)

    _rect(slide, MARGIN + col_w + Inches(0.28), cols_top, Pt(1.5), bottom - cols_top, tokens["muted"])

    _text(slide, data.get("right_label", ""), right_left, cols_top, col_w, Inches(0.4),
          17, tokens["accent"], tokens["font_title"], bold=True)
    _multi_para(slide, data.get("right_points") or [], right_left, cols_top + Inches(0.45), col_w, bottom - cols_top - Inches(0.45),
                14, tokens["text"], tokens["font_body"], marker="•  ", space_after=8)


def _render_causal_chain(slide, data, tokens, image_bytes):
    top = _content_top(tokens)
    _text(slide, data.get("process_title", data.get("title", "")), MARGIN, top, SLIDE_W - 2 * MARGIN, Inches(0.9),
          24, tokens["title"], tokens["font_title"], bold=True)

    row_top = top + Inches(1.1)
    box_w = Inches(3.5)
    gap = Inches(0.55)
    labels = [("Причина", data.get("cause", "")), ("Механизм", data.get("mechanism", "")), ("Итог", data.get("effect", ""))]
    x = MARGIN
    row_h = Inches(2.6) if not image_bytes else Inches(2.1)
    for i, (label, text) in enumerate(labels):
        card = _rect(slide, x, row_top, box_w, row_h, tokens.get("bg2") or tokens["bg"])
        if tokens["mode"] == "light":
            card.fill.fore_color.rgb = _hex("FFFFFF" if tokens["decoration"] != "rules" else tokens["bg"])
            card.line.fill.solid()
            card.line.color.rgb = _hex(tokens["muted"])
            card.line.width = Pt(0.75)
        _text(slide, label.upper(), x + Inches(0.2), row_top + Inches(0.15), box_w - Inches(0.4), Inches(0.35),
              12, tokens["accent"], tokens["font_title"], bold=True)
        _text(slide, text, x + Inches(0.2), row_top + Inches(0.55), box_w - Inches(0.4), row_h - Inches(0.75),
              13, tokens["text"], tokens["font_body"], line_spacing=1.2)
        x += box_w + gap
        if i < 2:
            _text(slide, "\u2192", x - gap + Inches(0.05), row_top + row_h / 2 - Inches(0.25), gap - Inches(0.1),
                  Inches(0.5), 22, tokens["accent"], tokens["font_title"], bold=True, align=PP_ALIGN.CENTER)

    if image_bytes:
        img_top = row_top + row_h + Inches(0.3)
        _image_or_placeholder(slide, image_bytes, MARGIN, img_top, SLIDE_W - 2 * MARGIN,
                               _content_bottom(tokens) - img_top, tokens)


def _render_quote_context(slide, data, tokens, image_bytes):
    top = _content_top(tokens)
    bottom = _content_bottom(tokens)
    img_w = Inches(4.6)
    _image_or_placeholder(slide, image_bytes, MARGIN, top, img_w, bottom - top, tokens)

    text_left = MARGIN + img_w + Inches(0.5)
    text_w = SLIDE_W - MARGIN - text_left

    context = data.get("context", "")
    quote = data.get("quote", "")
    explanation = data.get("explanation", "")

    context_h = Inches(0.8)
    _, context_used_h = _text(slide, context, text_left, top, text_w, context_h,
                               14, tokens["muted"], tokens["font_body"], italic=True)
    quote_top = top + Inches(max(context_used_h + 0.15, 0.6))
    quote_h = Inches(1.8)
    _, quote_used_h = _text(slide, f"\u00AB{quote}\u00BB", text_left, quote_top, text_w, quote_h,
                             22, tokens["title"], tokens["font_title"], bold=True, line_spacing=1.2)
    explanation_top = quote_top + Inches(max(quote_used_h + 0.2, 1.0))
    _text(slide, explanation, text_left, explanation_top, text_w, bottom - explanation_top,
          14, tokens["text"], tokens["font_body"], line_spacing=1.3)


def _render_conclusion(slide, data, tokens, image_bytes):
    _text(slide, data.get("title", ""), Inches(1.2), Inches(1.6), SLIDE_W - Inches(2.4), Inches(1.0),
          32, tokens["title"], tokens["font_title"], bold=True, align=PP_ALIGN.CENTER)
    bullets = data.get("bullets") or []
    _multi_para(slide, bullets, Inches(2.2), Inches(3.0), SLIDE_W - Inches(4.4), Inches(3.5),
                17, tokens["text"], tokens["font_body"], marker="\u2713  ", space_after=14, align=PP_ALIGN.LEFT)


LAYOUT_RENDERERS = {
    "title": _render_title,
    "toc": _render_toc,
    "thesis_proof": _render_thesis_proof,
    "dense_text": _render_dense_text,
    "concept_anatomy": _render_concept_anatomy,
    "comparison": _render_comparison,
    "causal_chain": _render_causal_chain,
    "quote_context": _render_quote_context,
    "conclusion": _render_conclusion,
}


async def build_pptx(title: str, slides: List[dict], style: str = "minimal", language: str = "ru") -> bytes:
    """
    slides: список dict с обязательным ключом "layout" и полями под эту раскладку
    (см. generators.LAYOUT_SCHEMAS), плюс "image_url" и "notes".
    """
    tokens = STYLE_TOKENS.get(style, STYLE_TOKENS["minimal"])

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank_layout = prs.slide_layouts[6]
    slide_count = len(slides)

    for idx, slide_data in enumerate(slides):
        slide = prs.slides.add_slide(blank_layout)
        layout = slide_data.get("layout", "thesis_proof")
        is_title_slide = layout == "title"

        draw_style_frame(slide, style, tokens, idx, slide_count, is_title_slide)

        image_bytes = None
        if slide_data.get("image_url"):
            image_bytes = await _download_image(slide_data["image_url"])

        renderer = LAYOUT_RENDERERS.get(layout, _render_thesis_proof)
        try:
            renderer(slide, slide_data, tokens, image_bytes)
        except Exception:
            logger.exception("Ошибка рендера слайда %d (layout=%s, style=%s) — safe fallback", idx, layout, style)
            _text(slide, slide_data.get("title", ""), MARGIN, Inches(0.9), SLIDE_W - 2 * MARGIN, Inches(1.0),
                  26, tokens["title"], tokens["font_title"], bold=True)
            fallback_lines = slide_data.get("facts") or slide_data.get("bullets") or slide_data.get("items") or []
            if fallback_lines:
                _multi_para(slide, fallback_lines, MARGIN, Inches(2.1), SLIDE_W - 2 * MARGIN, Inches(4.5),
                            16, tokens["text"], tokens["font_body"], marker="•  ")

        if slide_data.get("notes"):
            slide.notes_slide.notes_text_frame.text = slide_data["notes"]

    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()
